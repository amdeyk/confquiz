import os
import subprocess
from pathlib import Path
from typing import List, Optional
from fastapi import UploadFile
from PIL import Image
from pptx import Presentation
import uuid
from config import settings


class MediaService:
    def __init__(self):
        self.upload_dir = Path(settings.upload_dir)
        self.slides_dir = Path(settings.slides_dir)
        self.thumbs_dir = Path(settings.thumbs_dir)

        # Create directories
        self.upload_dir.mkdir(parents=True, exist_ok=True)
        self.slides_dir.mkdir(parents=True, exist_ok=True)
        self.thumbs_dir.mkdir(parents=True, exist_ok=True)

    async def save_ppt(self, file: UploadFile) -> str:
        """Save uploaded PPT file"""
        # Generate unique filename
        file_ext = Path(file.filename).suffix
        unique_filename = f"{uuid.uuid4()}{file_ext}"
        file_path = self.upload_dir / unique_filename

        # Save file
        with open(file_path, "wb") as f:
            content = await file.read()
            f.write(content)

        return str(file_path)

    def convert_ppt_to_images(self, ppt_path: str, deck_id: int) -> List[dict]:
        """Convert PPT to PNG images"""
        slides = []

        try:
            # Try LibreOffice conversion first (if available)
            if settings.libreoffice_path and os.path.exists(settings.libreoffice_path):
                slides = self._convert_with_libreoffice(ppt_path, deck_id)
            else:
                # Fallback to python-pptx (limited rendering)
                slides = self._convert_with_pptx(ppt_path, deck_id)

        except Exception as e:
            print(f"Error converting PPT: {e}")
            # Fallback method
            slides = self._convert_with_pptx(ppt_path, deck_id)

        return slides

    def _convert_with_libreoffice(self, ppt_path: str, deck_id: int) -> List[dict]:
        """Convert using LibreOffice (better quality)"""
        output_dir = self.slides_dir / f"deck_{deck_id}"
        output_dir.mkdir(parents=True, exist_ok=True)

        # Convert PPT to PDF first
        pdf_output = output_dir / "temp.pdf"

        # Run LibreOffice conversion
        cmd = [
            settings.libreoffice_path,
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', str(output_dir),
            ppt_path
        ]

        try:
            subprocess.run(cmd, check=True, timeout=60)

            # Convert PDF to images using PIL (requires pdf2image)
            from pdf2image import convert_from_path

            images = convert_from_path(str(pdf_output))
            slides = []

            for i, image in enumerate(images):
                slide_filename = f"slide_{i:03d}.png"
                thumb_filename = f"thumb_{i:03d}.png"

                slide_path = output_dir / slide_filename
                thumb_path = self.thumbs_dir / f"deck_{deck_id}_{thumb_filename}"

                # Save full size
                image.save(slide_path, "PNG")

                # Create thumbnail
                image.thumbnail((200, 150))
                image.save(thumb_path, "PNG")

                slides.append({
                    "slide_index": i,
                    "png_path": str(slide_path),
                    "thumb_path": str(thumb_path)
                })

            # Clean up PDF
            pdf_output.unlink(missing_ok=True)

            return slides

        except Exception as e:
            print(f"LibreOffice conversion failed: {e}")
            raise

    def _convert_with_pptx(self, ppt_path: str, deck_id: int) -> List[dict]:
        """Convert using python-pptx (basic rendering)"""
        output_dir = self.slides_dir / f"deck_{deck_id}"
        output_dir.mkdir(parents=True, exist_ok=True)

        # Load presentation
        prs = Presentation(ppt_path)
        slides = []

        for i, slide in enumerate(prs.slides):
            slide_filename = f"slide_{i:03d}.png"
            thumb_filename = f"thumb_{i:03d}.png"

            slide_path = output_dir / slide_filename
            thumb_path = self.thumbs_dir / f"deck_{deck_id}_{thumb_filename}"

            # Create blank image (since python-pptx doesn't render)
            # In production, you'd need proper rendering
            width = int(prs.slide_width.inches * 96)
            height = int(prs.slide_height.inches * 96)

            image = Image.new('RGB', (width, height), color='white')

            # Save full size
            image.save(slide_path, "PNG")

            # Create thumbnail
            image.thumbnail((200, 150))
            image.save(thumb_path, "PNG")

            slides.append({
                "slide_index": i,
                "png_path": str(slide_path),
                "thumb_path": str(thumb_path)
            })

        return slides

    def get_slide_count(self, ppt_path: str) -> int:
        """Get number of slides in presentation"""
        try:
            prs = Presentation(ppt_path)
            return len(prs.slides)
        except Exception as e:
            print(f"Error reading PPT: {e}")
            return 0


# Global instance
media_service = MediaService()
